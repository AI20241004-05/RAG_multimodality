# 第一支程式：建立多檔案的 Chroma 資料庫（PDF, PPTX, PPT），原始資料檔置於 C:\Users\user\Material\Report\source

from langchain.document_loaders import (PyMuPDFLoader, UnstructuredPowerPointLoader, 
     UnstructuredWordDocumentLoader, TextLoader, CSVLoader, UnstructuredMarkdownLoader,
     UnstructuredExcelLoader, UnstructuredHTMLLoader)
from langchain_unstructured import UnstructuredLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
#from langchain.chains import RetrievalQA
#from langchain_openai import ChatOpenAI
import os
import json
from langchain.schema import Document
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document as DocxDocument  # 處理 .docx 檔案
import io
import win32com.client
import datetime
#from langchain_community.vectorstores.utils import filter_complex_metadata

# 設定向量資料庫的儲存路徑
persist_directory = 'db'
model_name = "sentence-transformers/all-MiniLM-L6-v2"
model_kwargs = {'device': 'cuda'}  # 或 {'device': 'cpu'}，視硬體設備而定

# 初始化嵌入模型
embedding = HuggingFaceEmbeddings(model_name=model_name, model_kwargs=model_kwargs)

# OCR 讀取圖像文本
def extract_text_from_image(image_path):
    try:
        img = Image.open(image_path).convert("RGB")  # 確保是 RGB 格式
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        print(f"Error processing {image_path}: {e}")
        return ""

# 解析 PDF 內文字與圖片 OCR
def extract_text_from_pdf(pdf_path):
    text = ""
    doc = fitz.open(pdf_path)

    for page in doc:
        text += page.get_text("text") + "\n"

        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_path = f"temp_image_{img_index}.png"

            with open(image_path, "wb") as f:
                f.write(image_bytes)
            text += extract_text_from_image(image_path) + "\n"
            os.remove(image_path)
    return text.strip()

# 解析 PPT 內文字與圖片 OCR
def extract_text_from_ppt(ppt_path):
    text = ""
    prs = Presentation(ppt_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            if hasattr(shape, "image"):
                image_stream = shape.image.blob
                image_path = "temp_image.png"
                with open(image_path, "wb") as f:
                    f.write(image_stream)
                text += extract_text_from_image(image_path) + "\n"
                os.remove(image_path)
    return text.strip()

def convert_doc_to_docx(doc_path):
    if not doc_path.lower().endswith(".doc"):
        return doc_path  # 如果已經是 .docx，直接返回

    word = win32com.client.Dispatch("Word.Application")
    doc = word.Documents.Open(os.path.abspath(doc_path))
    
    new_path = doc_path[:len(doc_path)-4:]+"-reformated"+doc_path[len(doc_path)-4::]+ "x"  # 轉為 .docx

    doc.SaveAs(new_path, FileFormat=16)  # 16 表示 .docx 格式
    doc.Close()
    word.Quit()

    return new_path  # 返回新檔案路徑


def extract_text_from_docx(docx_path):
    text = ""
    doc = DocxDocument(docx_path)
    
    # 提取純文字
    for para in doc.paragraphs:
        text += para.text + "\n"

    # 提取圖片並進行 OCR
    for rel in doc.part.rels:
        if "image" in doc.part.rels[rel].target_ref:
            image_part = doc.part.rels[rel].target_part
            image_bytes = image_part.blob
            image = Image.open(io.BytesIO(image_bytes))
            text += pytesseract.image_to_string(image) + "\n"

    return text.strip()

# 處理多個檔案的載入函數
def load_documents(folder_path):
    documents = []
    for file in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file)
        if file.endswith(".pdf"):
            text = extract_text_from_pdf(file_path)
            if text:
                #documents.append({"text": text, "source": file_path})  
                documents.append(Document(page_content=text, metadata={"source": file_path}))  # 修正為 Document
            """
            loader = PyMuPDFLoader(file_path)
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
            """
        elif file.endswith(".pptx") or file.endswith(".ppt"):
            text = extract_text_from_ppt(file_path)
            if text:
                documents.append(Document(page_content=text, metadata={"source": file_path})) 
            """
            loader = UnstructuredPowerPointLoader(file_path)
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
            """    
        elif file.endswith(".docx") or file.endswith(".doc"):
            file_path = convert_doc_to_docx(file_path)  # 轉換為 .docx
            text = extract_text_from_docx(file_path)
            if text:
                documents.append(Document(page_content=text, metadata={"source": file_path})) 
            """
            loader = UnstructuredWordDocumentLoader(file_path)
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
            """    
        elif file.endswith(".txt"):
            loader = TextLoader(file_path)
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
        elif file.endswith(".csv"):
            loader = CSVLoader(file_path)
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
        elif file.endswith(".md"):
            loader = UnstructuredMarkdownLoader(file_path)
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
        elif file.endswith(".xlsx") or file.endswith(".xls"):
            loader = UnstructuredExcelLoader(file_path)
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
        elif file.endswith(".html"):
            loader = UnstructuredHTMLLoader(file_path)
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
        elif file.endswith((".js", ".py", ".dart", ".php", ".xml")):
            #loader = UnstructuredFileLoader(file_path)  # 使用通用載入器來處理這些格式
            loader = UnstructuredLoader(file_path)  # 使用更新的載入器
            loaded_docs = loader.load()
            # 將每個文件轉換為 Document 物件
            for doc in loaded_docs:
                documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
        elif file.endswith((".json", ".ipynb")):
            # JSON 需特別處理，不用 UnstructuredLoader
            with open(file_path, "r", encoding="utf-8") as f:
                json_data = json.load(f)
                # 轉換 JSON 內容成文本（可以依你的需求調整）
                json_text = json.dumps(json_data, indent=2, ensure_ascii=False)
                #documents.append(json_text)  # 直接存成文本
                # 將 JSON 文本轉換為 Document 物件
                documents.append(Document(page_content=json_text, metadata={"source": file_path}))
            continue  # 避免 UnstructuredLoader 再去處理 JSON
        elif file.endswith((".jpg", ".png")):
            text = extract_text_from_image(file_path)
            if text:
                #documents.append({"text": text, "source": file_path})
                # OCR 讀取後轉換為 Document
                documents.append(Document(page_content=text, metadata={"source": file_path}))
        else:
            continue
        #documents.extend(loader.load())
        # 將每個文件轉換為 Document 物件
        #loaded_docs = loader.load()
        #for doc in loaded_docs:
        #    documents.append(Document(page_content=doc.page_content, metadata=doc.metadata))
    return documents

# 分割文本
#text_splitter = RecursiveCharacterTextSplitter(chunk_size=100, chunk_overlap=10)
text_splitter = RecursiveCharacterTextSplitter(chunk_size=100, chunk_overlap=5)

# ct stores current time
ct = datetime.datetime.now()
print(f"Start time: {ct}")
# ts store timestamp of current time
ts = ct.timestamp()

# 指定要載入檔案的資料夾
folder_path = "C:\\Users\\user\\Material\\Report\\source"
all_documents = load_documents(folder_path)

# 分割所有文件
all_splits = text_splitter.split_documents(all_documents)

# 轉換 metadata 中的 list 為 string, 2/18
for doc in all_splits:
    if doc.metadata:
        for key, value in doc.metadata.items():
            if isinstance(value, list):  # 如果是 list，就轉成字串
                doc.metadata[key] = ", ".join(map(str, value))

# 儲存嵌入到 Chroma 資料庫
vectordb = Chroma.from_documents(documents=all_splits, embedding=embedding, persist_directory=persist_directory)

#vectordb.persist()
ct2 = datetime.datetime.now()
print(f"End time: {ct2}")

# ts store timestamp of current time
ts2 = ct2.timestamp()
print(f"{(ts2-ts)/60} min")
print("資料已成功儲存至 Chroma 資料庫")